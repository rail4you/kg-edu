#!/usr/bin/env python3
"""
generate_pptx.py — Generate a PPTX presentation from structured content.

Uses python-pptx library for professional slide generation.

Input (JSON via argv[1]):
{
  "courseName": "课程名称",
  "slides": "[{\"title\":\"Slide 1\",\"content\":\"...\",\"bullets\":[\"point 1\",\"point 2\"]}]",
  "author": "Author Name",
  "outputDir": "/tmp"
}

Output (JSON to stdout):
{"filePath": "/tmp/课程名称.pptx", "slideCount": 5}
"""
import json, sys, os, uuid
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(prs, slide_data, layout_idx=1):
    """Add a content slide with title and body."""
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = slide_data.get("title", "")

    # Body content
    content = slide_data.get("content", "")
    bullets = slide_data.get("bullets", [])

    # Find or add text body placeholder
    body_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body_shape = shape
            break

    if body_shape is None and len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]

    if body_shape:
        tf = body_shape.text_frame
        tf.clear()

        if content:
            p = tf.paragraphs[0]
            p.text = content
            p.level = 0

        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0


def main():
    try:
        data = json.loads(sys.argv[1])
        course_name = data["courseName"]
        slides_raw = data["slides"]
        author = data.get("author", "KgEdu Platform")
        output_dir = data["outputDir"]

        # Parse slides JSON (might be string or already parsed array)
        if isinstance(slides_raw, str):
            slides = json.loads(slides_raw)
        else:
            slides = slides_raw

        if not isinstance(slides, list) or len(slides) == 0:
            slides = [{"title": course_name, "content": "暂无内容"}]

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # --- Title slide ---
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = course_name
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"作者: {author}"

        # --- Content slides ---
        for s in slides:
            add_slide(prs, s)

        # Save
        safe_name = course_name.replace("/", "_").replace(" ", "_")[:50]
        file_name = f"{safe_name}_{uuid.uuid4().hex[:8]}.pptx"
        file_path = os.path.join(output_dir, file_name)
        prs.save(file_path)

        print(json.dumps({
            "filePath": file_path,
            "slideCount": len(prs.slides)
        }))

    except Exception as e:
        print(json.dumps({"error": str(e)}), file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
