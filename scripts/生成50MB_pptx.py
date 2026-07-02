"""
生成约 50MB 的 PPTX 文件
80 页，每页一张 1920x1080 渐变 + 文字图
"""
import os
import time
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.util import Inches


def get_font(size, bold=False):
    candidates = [
        '/System/Library/Fonts/PingFang.ttc',
        '/System/Library/Fonts/STHeiti Medium.ttc',
        '/System/Library/Fonts/STHeiti Light.ttc',
        '/System/Library/Fonts/Hiragino Sans GB.ttc',
        '/Library/Fonts/Arial Unicode.ttf',
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                continue
    return ImageFont.load_default()


FONT_BIG = get_font(150, True)
FONT_MID = get_font(90, True)
FONT_SM = get_font(50)


THEMES = [
    ('#16365D', '#2E5C8A', '#4A8FE7', '#A8D0F5'),
    ('#1A3C5E', '#2D5F8A', '#5BA0D0', '#B0DBF0'),
    ('#2D1B4E', '#5C2D7A', '#9B5BBF', '#D7B8E8'),
    ('#4A1942', '#7A2C5F', '#C45B95', '#F0B8D7'),
    ('#1B3A1B', '#2D5C2D', '#5BA05B', '#B0DFB0'),
    ('#3D2B1B', '#6B4A2D', '#B5825B', '#E5C9A8'),
    ('#3D1B1B', '#6B2D2D', '#B55B5B', '#E8B8B8'),
    ('#1B3D3D', '#2D6B6B', '#5BB5B5', '#B0E0E0'),
    ('#0D1B2A', '#1B3A5C', '#3D6FB0', '#88B5E0'),
    ('#2C1810', '#5C3A24', '#B07A4A', '#E5C49A'),
    ('#1A1A2E', '#3A3A5C', '#6B6B9B', '#B5B5D8'),
    ('#0F2A1A', '#2D5C3A', '#5BA07A', '#B0DFC0'),
    ('#3A1A3A', '#6B2D6B', '#B05BB0', '#E0B8E0'),
    ('#2A0F1A', '#5C2D3A', '#A05B7A', '#E5B8C9'),
    ('#1A2A3A', '#2D4A6B', '#5B7FA0', '#B0C9E0'),
]


TITLES = [
    '智慧教学 · 知识图谱', 'AI 智能体 · 自动出题', '思政教育 · 课程融合',
    '微专业 · 课程体系', '课堂互动 · 实时签到', '学情分析 · 数据洞察',
    '学习路径 · 进度跟踪', 'AI 问答 · 多轮对话', '智慧空间 · 智能工具',
    '教学资源 · 智能推荐', '考试管理 · 智能批改', '教师管理 · 多人协作',
    '学生管理 · 分组教学', '租户管理 · 独立部署', '多级组织 · 权限管理',
    '邮件问答 · 师生互通', '超星兼容 · 模板导入', '通义千问 · 私有化',
    '国产化 · 信创适配', '数据迁移 · 标准接口', '能力图谱 · 关系可视化',
    '思政图谱 · 案例库', '多轮对话 · 上下文', '教案生成 · AI 辅助',
    'PPT 生成 · 一键导出', 'DOCX 生成 · 自动排版', 'OSS 存储 · 数据安全',
    'API KEY · 动态更换', '活动日志 · 行为审计', '班级分析 · 整体学情',
    '知识点掌握 · 雷达图', '学习建议 · 个性化', '错题分析 · 弱项识别',
    '考试分析 · 成绩分布', '资源统计 · 使用频率', '学习推荐 · 智能匹配',
    '课程体系 · 学期规划', '章节管理 · 树形结构', '作业批改 · 智能评分',
    '视频学习 · 进度同步', '实验课程 · 虚拟仿真', '在线讨论 · 师生互动',
    '互动课堂 · 抢答抽奖', '学生端 · 学习工作台', '教师端 · 教学控制台',
    '管理端 · 租户配置', '超级管理员 · 跨租户', '院系管理 · 多级组织',
    '班级管理 · 学生分组', '成绩分析 · 多维度', '考试安排 · 时间规划',
    '题库管理 · 智能出卷', '试卷管理 · 难度分级', '主观题 · AI 批改',
    '客观题 · 自动评分', '学生答题 · 数据采集', '教师端 · AI 助教',
    '学情报告 · 导出', '教学反思 · 数据支撑', '教研协作 · 集体备课',
    '课程标准 · 知识映射', '岗位分析 · 能力对接', '专业能力 · 培养方案',
    '核心能力 · 雷达图', '通识能力 · 测评', '实践能力 · 项目制',
    '微专业发布 · 选课', '微专业学习 · 时间线', '微专业作业 · 提交批改',
    '期末评估 · 综合素养', '知识图谱 · 可视化', '思政元素 · 课程思政',
    '教学督导 · 数据驾驶舱', '教研活动 · 记录', '教师培训 · AI 助手',
    '学生画像 · 综合分析', '就业指导 · 能力匹配', '创新创业 · 实践平台',
    '产教融合 · 案例库', '在线实验 · 资源池', '虚拟仿真 · 沉浸式',
    '毕业设计 · 全流程', '论文管理 · 查重',
]


def hex2rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def make_gradient(w, h, c1, c2, angle=90):
    img = Image.new('RGB', (w, h), c1)
    draw = ImageDraw.Draw(img)
    r1, g1, b1 = c1
    r2, g2, b2 = c2
    if angle == 90:
        for y in range(h):
            t = y / h
            draw.line([(0, y), (w, y)], fill=(int(r1 + (r2-r1)*t), int(g1 + (g2-g1)*t), int(b1 + (b2-b1)*t)))
    else:
        step = 4
        for y in range(0, h, step):
            for x in range(0, w, step):
                t = (x + y) / (w + h)
                draw.rectangle([(x, y), (x+step, y+step)],
                               fill=(int(r1 + (r2-r1)*t), int(g1 + (g2-g1)*t), int(b1 + (b2-b1)*t)))
    return img


def add_decorations(img, c3, c4):
    w, h = img.size
    overlay = Image.new('RGBA', (w, h), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    r = int(min(w, h) * 0.35)
    od.ellipse([(-r//2, -r//2), (r, r)], fill=c3 + (60,))
    overlay = overlay.filter(ImageFilter.GaussianBlur(80))
    img.paste(overlay, (0, 0), overlay)

    overlay2 = Image.new('RGBA', (w, h), (0, 0, 0, 0))
    od2 = ImageDraw.Draw(overlay2)
    r2 = int(min(w, h) * 0.25)
    od2.ellipse([(w - r2//2, h - r2//2), (w + r2, h + r2)], fill=c4 + (80,))
    overlay2 = overlay2.filter(ImageFilter.GaussianBlur(60))
    img.paste(overlay2, (0, 0), overlay2)

    draw = ImageDraw.Draw(img, 'RGBA')
    for i in range(8):
        y = h * 0.2 + i * 40
        draw.line([(0, y), (w, y)], fill=(255, 255, 255, 12), width=2)
    return img


def add_text(img, title, page_num, total, theme_idx):
    draw = ImageDraw.Draw(img, 'RGBA')
    w, h = img.size
    # 顶部品牌条（按宽度比例）
    bar_h = int(h * 0.055)
    draw.rectangle([(0, 0), (w, bar_h)], fill=(0, 0, 0, 80))
    draw.text((int(w * 0.018), int(bar_h * 0.2)), '课堂星 · 智慧教学系统', font=FONT_SM, fill=(255, 255, 255, 220))

    bbox = draw.textbbox((0, 0), title, font=FONT_BIG)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    x = (w - tw) // 2
    y = h // 2 - th
    for dx, dy in [(6, 6), (-4, 4), (4, -4)]:
        draw.text((x + dx, y + dy), title, font=FONT_BIG, fill=(0, 0, 0, 100))
    draw.text((x, y), title, font=FONT_BIG, fill=(255, 255, 255, 255))

    sub = THEMES[theme_idx][3].upper() + ' · ' + str(theme_idx + 1).zfill(2)
    bbox2 = draw.textbbox((0, 0), sub, font=FONT_MID)
    sw_ = bbox2[2] - bbox2[0]
    draw.text(((w - sw_) // 2, y + th + 30), sub, font=FONT_MID, fill=(255, 255, 255, 200))

    # 底部页码 + 主题（按比例定位）
    page_text = f'{page_num:02d} / {total:02d}'
    draw.text((w - int(w * 0.10), h - int(h * 0.10)), page_text, font=FONT_MID, fill=(255, 255, 255, 180))
    draw.text((int(w * 0.02), h - int(h * 0.10)), f'THEME {theme_idx + 1:02d}', font=FONT_MID, fill=(255, 255, 255, 180))
    return img


def make_slide_image(page_num, total, theme_idx, title):
    # 3000x1688 单图 ~493KB，100 张目标 50MB
    w, h = 3000, 1688
    c1, c2, c3, c4 = [hex2rgb(c) for c in THEMES[theme_idx]]
    img = make_gradient(w, h, c1, c2, angle=90 if page_num % 2 == 0 else 45)
    img = add_decorations(img, c3, c4)
    img = add_text(img, title, page_num, total, theme_idx)
    img = img.filter(ImageFilter.SMOOTH)
    return img


def main():
    total_slides = 200
    out_path = '/Users/bai/projects/kg-edu/docs/课堂星智慧教学系统_50MB.pptx'
    tmp_dir = '/tmp/pptx_imgs_50mb'
    os.makedirs(tmp_dir, exist_ok=True)

    print(f'开始生成 {total_slides} 张 1920x1080 渐变图...')
    start = time.time()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for i in range(total_slides):
        title = TITLES[i] if i < len(TITLES) else f'功能展示 {i+1:02d}'
        theme_idx = i % len(THEMES)
        img = make_slide_image(i + 1, total_slides, theme_idx, title)

        tmp_path = os.path.join(tmp_dir, f'slide_{i:03d}.jpg')
        # quality=100 + 无子采样，单图 ~500KB，目标 50MB
        img.save(tmp_path, 'JPEG', quality=100, optimize=False, subsampling=0)

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(tmp_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        if (i + 1) % 10 == 0:
            elapsed = time.time() - start
            print(f'  进度: {i+1}/{total_slides} ({elapsed:.1f}s)')

    print('保存 PPTX...')
    prs.save(out_path)
    size_mb = os.path.getsize(out_path) / (1024 * 1024)
    print(f'\n✅ 完成: {out_path}')
    print(f'   文件大小: {size_mb:.2f} MB')
    print(f'   总耗时: {time.time()-start:.1f}s')

    for f in os.listdir(tmp_dir):
        os.remove(os.path.join(tmp_dir, f))
    os.rmdir(tmp_dir)


if __name__ == '__main__':
    main()
